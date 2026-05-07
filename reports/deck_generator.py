"""
PhantomFeed — PowerPoint Briefing Deck Generator
Generates a 12-slide executive threat intelligence briefing.
"""
import io
import json
import os
from datetime import datetime, timedelta
from typing import Optional

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    import matplotlib.patches as mpatches
    MATPLOTLIB_AVAILABLE = True
except ImportError:
    MATPLOTLIB_AVAILABLE = False

# PhantomFeed brand colours
NAVY   = RGBColor(0x1a, 0x20, 0x35) if PPTX_AVAILABLE else None
TEAL   = RGBColor(0x00, 0xd4, 0xaa) if PPTX_AVAILABLE else None
WHITE  = RGBColor(0xFF, 0xFF, 0xFF) if PPTX_AVAILABLE else None
GRAY   = RGBColor(0x64, 0x74, 0x8b) if PPTX_AVAILABLE else None
RED    = RGBColor(0xFF, 0x47, 0x57) if PPTX_AVAILABLE else None
AMBER  = RGBColor(0xFF, 0xa5, 0x02) if PPTX_AVAILABLE else None
DARKBG = RGBColor(0x0d, 0x0f, 0x14) if PPTX_AVAILABLE else None

SEV_COLORS = {
    "CRITICAL": (0xFF, 0x47, 0x57),
    "HIGH":     (0xFF, 0xa5, 0x02),
    "MEDIUM":   (0x00, 0xb4, 0xd8),
    "LOW":      (0x64, 0x74, 0x8b),
}


def _rgb(r, g, b):
    return RGBColor(r, g, b) if PPTX_AVAILABLE else None


def _slide_bg(slide, prs, color=None):
    """Set slide background to dark navy."""
    if not PPTX_AVAILABLE:
        return
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color or NAVY


def _text_box(slide, left, top, width, height, text, font_size=12,
              bold=False, color=None, align=PP_ALIGN.LEFT if PPTX_AVAILABLE else None,
              wrap=True):
    if not PPTX_AVAILABLE:
        return None
    from pptx.util import Inches, Pt
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color or WHITE
    return txBox


def _header_bar(slide, title, subtitle="", prs=None):
    """Add a teal header bar at the top of the slide."""
    if not PPTX_AVAILABLE:
        return
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0), Inches(0), Inches(10), Inches(1.2)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()

    _text_box(slide, 0.2, 0.1, 9.6, 0.65, title,
              font_size=24, bold=True, color=DARKBG)
    if subtitle:
        _text_box(slide, 0.2, 0.75, 9.6, 0.4, subtitle,
                  font_size=10, color=NAVY)


def _make_bar_chart(labels, values, title, color="#00d4aa") -> Optional[io.BytesIO]:
    """Generate a horizontal bar chart and return PNG bytes."""
    if not MATPLOTLIB_AVAILABLE:
        return None
    fig, ax = plt.subplots(figsize=(8, max(3, len(labels) * 0.4)))
    fig.patch.set_facecolor("#1a1e27")
    ax.set_facecolor("#1a1e27")
    bars = ax.barh(labels, values, color=color, edgecolor="none")
    ax.set_xlabel("Count", color="#64748b")
    ax.set_title(title, color="#e2e8f0", pad=10)
    ax.tick_params(colors="#64748b")
    for spine in ax.spines.values():
        spine.set_edgecolor("#2a2f3e")
    for bar, val in zip(bars, values):
        ax.text(bar.get_width() + 0.1, bar.get_y() + bar.get_height() / 2,
                str(val), va="center", color="#e2e8f0", fontsize=9)
    plt.tight_layout()
    buf = io.BytesIO()
    plt.savefig(buf, format="png", facecolor=fig.get_facecolor(), dpi=100)
    plt.close()
    buf.seek(0)
    return buf


def _make_donut(data: dict, title: str) -> Optional[io.BytesIO]:
    """Generate a donut chart. data = {label: count}."""
    if not MATPLOTLIB_AVAILABLE:
        return None
    colors = ["#ff4757", "#ffa502", "#00d4aa", "#64748b"]
    labels = list(data.keys())
    values = [max(0, v) for v in data.values()]
    if sum(values) == 0:
        values = [1] * len(values)

    fig, ax = plt.subplots(figsize=(5, 5))
    fig.patch.set_facecolor("#1a1e27")
    ax.set_facecolor("#1a1e27")
    wedges, texts, autotexts = ax.pie(
        values, labels=labels, autopct="%1.0f%%",
        colors=colors[:len(labels)], startangle=90,
        wedgeprops={"width": 0.5, "edgecolor": "#0d0f14"},
        textprops={"color": "#e2e8f0", "fontsize": 9},
    )
    for at in autotexts:
        at.set_color("#0d0f14")
        at.set_fontweight("bold")
    ax.set_title(title, color="#e2e8f0", pad=10)
    plt.tight_layout()
    buf = io.BytesIO()
    plt.savefig(buf, format="png", facecolor=fig.get_facecolor(), dpi=100)
    plt.close()
    buf.seek(0)
    return buf


async def _ollama_summary(prompt: str, model: str = "llama3.2") -> str:
    """Call local Ollama for AI-generated content."""
    try:
        import httpx
        async with httpx.AsyncClient(timeout=60) as client:
            resp = await client.post(
                "http://localhost:11434/api/generate",
                json={"model": model, "prompt": prompt, "stream": False},
            )
            if resp.status_code == 200:
                return resp.json().get("response", "").strip()
    except Exception:
        pass
    return ""


class BriefingDeckGenerator:

    async def generate_deck(self, client_id: str, days: int = 30) -> bytes:
        """Generate a complete 12-slide PowerPoint briefing deck."""
        if not PPTX_AVAILABLE:
            raise RuntimeError("python-pptx not installed. Run: pip install python-pptx")

        from db import database as db
        client = await db.get_client(client_id)
        if not client:
            raise ValueError(f"Client {client_id} not found")

        stack = client.get("stack_profile") or {}
        industry = client.get("industry") or stack.get("industry") or "Technology"
        client_name = client.get("name", "Client")
        date_from = (datetime.utcnow() - timedelta(days=days)).strftime("%B %d, %Y")
        date_to = datetime.utcnow().strftime("%B %d, %Y")

        # Fetch data
        all_items = await db.get_items(limit=500, sort="risk", client_id=client_id)
        if not all_items:
            all_items = await db.get_items(limit=200, sort="risk")
        recent = [i for i in all_items if i.get("published_at", "") >= (datetime.utcnow() - timedelta(days=days)).strftime("%Y-%m-%d")]

        crit_items = [i for i in recent if i.get("severity") == "CRITICAL"]
        high_items = [i for i in recent if i.get("severity") == "HIGH"]
        kev_items  = [i for i in recent if i.get("category") == "kev"]

        remediations = await db.get_remediations(client_id)
        rem_stats = {
            "Open": sum(1 for r in remediations if r.get("status") == "open"),
            "In Progress": sum(1 for r in remediations if r.get("status") == "in_progress"),
            "Patched": sum(1 for r in remediations if r.get("status") == "patched"),
            "Overdue": sum(1 for r in remediations if r.get("is_overdue")),
        }

        # Vendor breakdown
        vendors: dict = {}
        for item in recent:
            v = item.get("vendor") or item.get("feed_label", "Other")
            vendors[v[:30]] = vendors.get(v[:30], 0) + 1
        top_vendors = sorted(vendors.items(), key=lambda x: x[1], reverse=True)[:8]

        # Compliance breakdown
        cmmc_hits: dict = {}
        for item in recent:
            for tag in (item.get("compliance_tags") or []):
                cmmc_hits[tag] = cmmc_hits.get(tag, 0) + 1
        top_cmmc = sorted(cmmc_hits.items(), key=lambda x: x[1], reverse=True)[:6]

        # AI-generated summaries
        threat_summary_prompt = (
            f"You are a cybersecurity analyst. Write 4 concise bullet points (each 1 sentence) "
            f"summarizing the threat landscape for a {industry} company over the last {days} days. "
            f"There were {len(crit_items)} CRITICAL, {len(high_items)} HIGH vulnerabilities, "
            f"and {len(kev_items)} CISA KEV items. "
            f"Top affected vendors: {', '.join(v for v, _ in top_vendors[:4])}. "
            f"Be professional and board-appropriate. No markdown, use plain bullet points with '-'."
        )
        threat_summary = await _ollama_summary(threat_summary_prompt)
        if not threat_summary:
            threat_summary = (
                f"- {len(crit_items)} critical vulnerabilities identified in the past {days} days\n"
                f"- {len(kev_items)} items appear on the CISA Known Exploited Vulnerabilities catalog\n"
                f"- {len(high_items)} high-severity advisories require prompt attention\n"
                f"- Immediate remediation recommended for CISA KEV items"
            )

        actions_prompt = (
            f"You are a CISO advising a {industry} company. "
            f"List 5 prioritized remediation actions for the next 30 days given: "
            f"{len(crit_items)} critical CVEs, {len(kev_items)} KEV items, "
            f"{rem_stats['Overdue']} overdue patches. "
            f"Be specific and actionable. No markdown, use '1.' '2.' etc."
        )
        actions_text = await _ollama_summary(actions_prompt)
        if not actions_text:
            actions_text = (
                "1. Immediately patch all CISA KEV items (legally binding for federal contractors)\n"
                f"2. Prioritize {len(crit_items)} CRITICAL severity vulnerabilities for emergency patching\n"
                "3. Review and update asset inventory to ensure complete coverage\n"
                "4. Verify EDR coverage across all endpoints\n"
                "5. Schedule tabletop exercise to test incident response procedures"
            )

        # Build presentation
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        blank_layout = prs.slide_layouts[6]

        def add_slide():
            sl = prs.slides.add_slide(blank_layout)
            _slide_bg(sl, prs)
            return sl

        # ── Slide 1: Title ──────────────────────────────────────────────────────
        sl = add_slide()
        # Brand bar at bottom
        br = sl.shapes.add_shape(1, Inches(0), Inches(6.8), Inches(10), Inches(0.7))
        br.fill.solid(); br.fill.fore_color.rgb = TEAL; br.line.fill.background()
        _text_box(sl, 0.2, 6.85, 5, 0.5, "PhantomFeed Threat Intelligence", 9, color=DARKBG)
        _text_box(sl, 7, 6.85, 3, 0.5, datetime.utcnow().strftime("%B %Y"), 9, color=DARKBG, align=PP_ALIGN.RIGHT)
        # Main title
        _text_box(sl, 0.5, 1.5, 9, 1.0, "Threat Intelligence Brief", 36, bold=True, color=TEAL)
        _text_box(sl, 0.5, 2.6, 9, 0.7, client_name, 28, bold=True)
        _text_box(sl, 0.5, 3.4, 9, 0.5, f"{date_from} — {date_to}", 14, color=GRAY)
        _text_box(sl, 0.5, 4.0, 9, 0.5, f"Industry: {industry} | Confidential", 11, color=GRAY)

        # ── Slide 2: Executive Summary ──────────────────────────────────────────
        sl = add_slide()
        _header_bar(sl, "Executive Summary", f"Threat Landscape | {date_from} – {date_to}")
        bullets = [b.strip().lstrip("- ") for b in threat_summary.split("\n") if b.strip()]
        y = 1.4
        for bullet in bullets[:6]:
            bar = sl.shapes.add_shape(1, Inches(0.4), Inches(y + 0.12), Inches(0.05), Inches(0.2))
            bar.fill.solid(); bar.fill.fore_color.rgb = TEAL; bar.line.fill.background()
            _text_box(sl, 0.6, y, 9.2, 0.5, bullet, 13)
            y += 0.65

        # ── Slide 3: Threat Landscape Overview ─────────────────────────────────
        sl = add_slide()
        _header_bar(sl, "Threat Landscape Overview", f"Summary Statistics | {days}-Day Period")
        stats = [
            ("CRITICAL", str(len(crit_items)), "Active Vulnerabilities", RED),
            ("HIGH", str(len(high_items)), "High Severity Items", AMBER),
            ("KEV", str(len(kev_items)), "Known Exploited", TEAL),
            ("TOTAL", str(len(recent)), "Items This Period", GRAY),
        ]
        x_positions = [0.3, 2.8, 5.3, 7.8]
        for (label, value, desc, color), x in zip(stats, x_positions):
            box = sl.shapes.add_shape(1, Inches(x), Inches(1.5), Inches(2.2), Inches(2.5))
            box.fill.solid(); box.fill.fore_color.rgb = _rgb(0x1a, 0x1e, 0x27)
            box.line.color.rgb = color
            _text_box(sl, x + 0.1, 1.6, 2, 0.4, label, 9, color=color, align=PP_ALIGN.CENTER)
            _text_box(sl, x + 0.1, 2.0, 2, 0.8, value, 36, bold=True, color=color, align=PP_ALIGN.CENTER)
            _text_box(sl, x + 0.1, 2.9, 2, 0.5, desc, 9, color=GRAY, align=PP_ALIGN.CENTER)

        # ── Slide 4: Top 5 Critical Threats ────────────────────────────────────
        sl = add_slide()
        _header_bar(sl, "Top 5 Critical Threats", "Highest Risk Items Requiring Immediate Action")
        top5 = (crit_items + high_items)[:5]
        y = 1.4
        for i, item in enumerate(top5):
            sev = item.get("severity", "HIGH")
            sc = SEV_COLORS.get(sev, (100, 100, 100))
            dot = sl.shapes.add_shape(1, Inches(0.3), Inches(y + 0.1), Inches(0.15), Inches(0.15))
            dot.fill.solid(); dot.fill.fore_color.rgb = _rgb(*sc); dot.line.fill.background()
            cves = ", ".join((item.get("cve_ids") or [])[:2])
            title_text = item.get("title", "")[:80]
            _text_box(sl, 0.6, y, 7, 0.35, title_text, 11, bold=True)
            meta = f"{sev}  |  {item.get('vendor', '')}  |  CVSS: {item.get('cvss', 'N/A')}  |  {cves}"
            _text_box(sl, 0.6, y + 0.35, 7, 0.3, meta, 9, color=GRAY)
            y += 0.85

        # ── Slide 5: CISA KEV Highlights ───────────────────────────────────────
        sl = add_slide()
        _header_bar(sl, "CISA KEV Highlights", "Known Exploited Vulnerabilities — Immediate Action Required")
        _text_box(sl, 0.3, 1.3, 9.5, 0.4,
                  "CISA KEV items are actively exploited in the wild. Federal contractors must patch within mandated timeframes.",
                  10, color=AMBER)
        kevs = kev_items[:6]
        y = 1.85
        for item in kevs:
            cves = ", ".join((item.get("cve_ids") or [])[:2])
            _text_box(sl, 0.3, y, 8, 0.28, f"- {item.get('title', '')[:90]} [{cves}]", 10)
            y += 0.42
        if not kevs:
            _text_box(sl, 0.3, 2.0, 9.5, 0.5, "No KEV items identified in this period.", 12, color=GRAY)

        # ── Slide 6: Threat Actor Activity ─────────────────────────────────────
        sl = add_slide()
        _header_bar(sl, "Threat Actor Activity", f"Active Groups Targeting {industry} Organizations")
        from db import database as db_ref
        actors = await db_ref.get_threat_actors(active_only=True)
        ind_low = industry.lower()
        targeting = [a for a in actors if any(
            ind_low in i.lower() or i.lower() in ind_low
            for i in (a.get("target_industries") or [])
        )][:5]
        y = 1.4
        for actor in targeting:
            ttps_str = ", ".join((actor.get("ttps") or [])[:4])
            _text_box(sl, 0.3, y, 3.5, 0.3, actor["name"], 12, bold=True, color=TEAL)
            _text_box(sl, 3.9, y, 2, 0.3, f"{actor.get('origin', '')} | {actor.get('motivation', '')}", 9, color=GRAY)
            _text_box(sl, 0.3, y + 0.32, 9, 0.3, f"TTPs: {ttps_str}", 9, color=GRAY)
            y += 0.8
        if not targeting:
            _text_box(sl, 0.3, 2.0, 9.5, 1.0,
                      f"No active threat actors with confirmed {industry} industry targeting in database.\nReview actors.html for full dossier browser.",
                      11, color=GRAY)

        # ── Slide 7: Vendor Risk Breakdown ─────────────────────────────────────
        sl = add_slide()
        _header_bar(sl, "Vendor Risk Breakdown", "Vulnerabilities by Affected Vendor")
        if top_vendors and MATPLOTLIB_AVAILABLE:
            chart_buf = _make_bar_chart(
                [v[:25] for v, _ in top_vendors],
                [c for _, c in top_vendors],
                "Items by Vendor",
                "#00d4aa",
            )
            if chart_buf:
                sl.shapes.add_picture(chart_buf, Inches(0.3), Inches(1.3), Inches(9.4), Inches(5.5))
        else:
            y = 1.5
            for vendor, count in top_vendors[:8]:
                _text_box(sl, 0.3, y, 7, 0.35, vendor, 12)
                _text_box(sl, 7.5, y, 2, 0.35, str(count), 12, bold=True, color=TEAL, align=PP_ALIGN.RIGHT)
                y += 0.5

        # ── Slide 8: Compliance Impact ──────────────────────────────────────────
        sl = add_slide()
        _header_bar(sl, "Compliance Impact", "CMMC 2.0 & NIST 800-171 Controls at Risk")
        _text_box(sl, 0.3, 1.3, 9.5, 0.35, "Controls affected by active vulnerabilities this period:", 11, color=GRAY)
        y = 1.75
        for tag, count in top_cmmc:
            domain = tag.split("-")[0]
            control = tag.split("-")[1] if "-" in tag else tag
            _text_box(sl, 0.3, y, 1.2, 0.35, domain, 10, bold=True, color=AMBER)
            _text_box(sl, 1.6, y, 1, 0.35, control, 10)
            _text_box(sl, 2.7, y, 5, 0.35, f"{count} item(s) affecting this control", 10, color=GRAY)
            y += 0.45
        if not top_cmmc:
            _text_box(sl, 0.3, 2.0, 9.5, 0.5, "No compliance tags found — ensure items are properly tagged.", 11, color=GRAY)

        # ── Slide 9: Asset Exposure Summary ────────────────────────────────────
        sl = add_slide()
        _header_bar(sl, "Asset Exposure Summary", "Confirmed Vulnerability Exposures")
        assets = await db_ref.get_assets(client_id)
        exposure_count = 0
        asset_types: dict = {}
        for asset in assets:
            t = asset.get("asset_type", "workstation")
            asset_types[t] = asset_types.get(t, 0) + 1
        _text_box(sl, 0.3, 1.4, 9.5, 0.4,
                  f"Total tracked assets: {len(assets)}  |  Asset types: {', '.join(f'{t} ({c})' for t, c in asset_types.items())}",
                  11, color=GRAY)
        if assets:
            y = 1.95
            for asset in assets[:8]:
                _text_box(sl, 0.3, y, 4, 0.3, asset.get("hostname") or asset.get("ip_address") or "Unknown", 10)
                _text_box(sl, 4.4, y, 3, 0.3, asset.get("software", ""), 10, color=GRAY)
                _text_box(sl, 7.5, y, 2, 0.3, asset.get("os", ""), 10, color=GRAY)
                y += 0.42
        else:
            _text_box(sl, 0.3, 2.0, 9.5, 0.5, "No assets configured — add assets via the Admin panel.", 11, color=GRAY)

        # ── Slide 10: Remediation Status ────────────────────────────────────────
        sl = add_slide()
        _header_bar(sl, "Remediation Status", "Open Items & Patch Velocity")
        if MATPLOTLIB_AVAILABLE and sum(rem_stats.values()) > 0:
            chart_buf = _make_donut(rem_stats, "Remediation Status")
            if chart_buf:
                sl.shapes.add_picture(chart_buf, Inches(0.3), Inches(1.3), Inches(4.5), Inches(4.5))
        # Stats on right
        x = 5.0; y = 1.5
        for label, val in rem_stats.items():
            color = RED if label == "Overdue" else TEAL if label == "Patched" else WHITE
            _text_box(sl, x, y, 2, 0.4, label, 12, color=GRAY)
            _text_box(sl, x + 2.2, y, 1.5, 0.4, str(val), 14, bold=True, color=color)
            y += 0.6
        if rem_stats.get("Overdue", 0) > 0:
            _text_box(sl, 5.0, y + 0.3, 4.5, 0.5,
                      f"ATTENTION: {rem_stats['Overdue']} items are past SLA deadline", 11, bold=True, color=RED)

        # ── Slide 11: Recommended Actions ───────────────────────────────────────
        sl = add_slide()
        _header_bar(sl, "Recommended Actions", f"Priority Remediation Plan — Next 30 Days")
        lines = [l.strip() for l in actions_text.split("\n") if l.strip()]
        y = 1.4
        for i, line in enumerate(lines[:7]):
            _text_box(sl, 0.3, y, 9.4, 0.5, line, 12)
            y += 0.65

        # ── Slide 12: Contact / Next Steps ──────────────────────────────────────
        sl = add_slide()
        _header_bar(sl, "Next Steps & Contact", "PhantomFeed Threat Intelligence Platform")
        import config as cfg
        _text_box(sl, 0.5, 1.5, 9, 0.4, "Your Dedicated Security Team", 14, bold=True, color=TEAL)
        contacts = [
            ("Provider", cfg.MSP_NAME),
            ("Email", cfg.MSP_EMAIL or "security@your-msp.com"),
            ("Phone", cfg.MSP_PHONE or "+1 (800) 000-0000"),
            ("Platform", "PhantomFeed Intelligence Dashboard"),
        ]
        y = 2.1
        for label, val in contacts:
            _text_box(sl, 0.5, y, 2, 0.35, label + ":", 11, color=GRAY)
            _text_box(sl, 2.6, y, 7, 0.35, val, 11, bold=True)
            y += 0.5
        _text_box(sl, 0.5, 4.5, 9, 0.5,
                  "Frequency: Weekly | Classification: Client Confidential | Generated by PhantomFeed",
                  9, color=GRAY)

        # Save to bytes
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return buf.read()

    async def get_deck_preview(self, client_id: str, days: int = 30) -> dict:
        """Return slide titles and content outline without generating the file."""
        from db import database as db
        client = await db.get_client(client_id)
        name = client.get("name", "Client") if client else "Client"
        items = await db.get_items(limit=200, sort="risk", client_id=client_id)
        recent = [i for i in items if i.get("published_at", "") >= (datetime.utcnow() - timedelta(days=days)).strftime("%Y-%m-%d")]
        crit = sum(1 for i in recent if i.get("severity") == "CRITICAL")
        high = sum(1 for i in recent if i.get("severity") == "HIGH")
        kev  = sum(1 for i in recent if i.get("category") == "kev")

        return {
            "client": name,
            "days": days,
            "slides": [
                {"slide": 1, "title": "Title Slide", "content": f"{name} — Threat Intelligence Brief"},
                {"slide": 2, "title": "Executive Summary", "content": f"AI-generated 4-bullet threat landscape summary"},
                {"slide": 3, "title": "Threat Landscape Overview", "content": f"{crit} CRITICAL, {high} HIGH, {kev} KEV items"},
                {"slide": 4, "title": "Top 5 Critical Threats", "content": "Highest risk items with CVSS scores"},
                {"slide": 5, "title": "CISA KEV Highlights", "content": f"{kev} actively exploited vulnerabilities"},
                {"slide": 6, "title": "Threat Actor Activity", "content": "Active APTs targeting client industry"},
                {"slide": 7, "title": "Vendor Risk Breakdown", "content": "Bar chart of items by affected vendor"},
                {"slide": 8, "title": "Compliance Impact", "content": "CMMC 2.0 / NIST controls at risk"},
                {"slide": 9, "title": "Asset Exposure Summary", "content": "Confirmed exposures by asset type"},
                {"slide": 10, "title": "Remediation Status", "content": "Donut chart: open/in-progress/patched/overdue"},
                {"slide": 11, "title": "Recommended Actions", "content": "AI-generated top 5 priority actions"},
                {"slide": 12, "title": "Contact / Next Steps", "content": "MSP contact information"},
            ],
        }

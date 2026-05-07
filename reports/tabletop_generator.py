"""
PhantomFeed — Tabletop Exercise Generator

Uses Ollama (local LLM) to generate scenario-specific tabletop exercises.
Exports to PDF (reportlab) or PPTX (python-pptx).
"""
import io
import json
from datetime import datetime
from typing import Optional


SCENARIO_TEMPLATES = {
    "ransomware": {
        "title": "Ransomware Attack Response",
        "phases": ["Initial Infection", "Lateral Movement", "Encryption Event", "Ransom Demand", "Recovery"],
        "prompt_hint": "ransomware attack that encrypts corporate file shares and demands cryptocurrency payment",
    },
    "supply_chain": {
        "title": "Software Supply Chain Compromise",
        "phases": ["Compromise Detection", "Scope Assessment", "Vendor Coordination", "Remediation", "Post-Incident Review"],
        "prompt_hint": "software supply chain compromise where a trusted vendor's update contains malicious code",
    },
    "data_breach": {
        "title": "Data Breach & Exfiltration",
        "phases": ["Breach Discovery", "Initial Triage", "Containment", "Notification", "Regulatory Response"],
        "prompt_hint": "data breach where customer PII is exfiltrated and appears on the dark web",
    },
    "insider_threat": {
        "title": "Malicious Insider Incident",
        "phases": ["Detection", "Investigation", "Containment", "HR/Legal Coordination", "Remediation"],
        "prompt_hint": "malicious insider threat where a disgruntled employee exfiltrates intellectual property",
    },
    "ddos": {
        "title": "Distributed Denial of Service Attack",
        "phases": ["Attack Onset", "Triage", "Mitigation", "Communication", "Recovery"],
        "prompt_hint": "volumetric DDoS attack targeting the organization's public-facing web services",
    },
    "phishing": {
        "title": "Spear Phishing & Credential Theft",
        "phases": ["Email Receipt", "Credential Compromise", "Account Takeover", "Lateral Movement", "Containment"],
        "prompt_hint": "targeted spear phishing campaign that results in executive email account compromise",
    },
    "zero_day": {
        "title": "Zero-Day Vulnerability Exploitation",
        "phases": ["Vulnerability Disclosure", "Exposure Assessment", "Emergency Patching", "Threat Hunting", "Lessons Learned"],
        "prompt_hint": "critical zero-day vulnerability in a widely used enterprise product actively exploited in the wild",
    },
    "cloud_breach": {
        "title": "Cloud Infrastructure Compromise",
        "phases": ["Initial Detection", "Cloud Resource Audit", "Credential Rotation", "Workload Isolation", "Recovery"],
        "prompt_hint": "cloud infrastructure compromise through misconfigured S3 buckets and stolen IAM credentials",
    },
}


async def _ollama_generate(prompt: str, model: str = "llama3.2") -> str:
    try:
        import httpx
        async with httpx.AsyncClient(timeout=120) as client:
            resp = await client.post(
                "http://localhost:11434/api/generate",
                json={"model": model, "prompt": prompt, "stream": False},
            )
            if resp.status_code == 200:
                return resp.json().get("response", "").strip()
    except Exception:
        pass
    return ""


def _fallback_scenario(template: dict, client_name: str, industry: str) -> dict:
    """Generate a deterministic fallback when Ollama is unavailable."""
    phases = []
    for phase in template["phases"]:
        phases.append({
            "phase": phase,
            "situation": f"The {phase.lower()} phase of the {template['title']} is underway.",
            "inject": f"[INJECT] {phase}: New information arrives requiring immediate decision.",
            "discussion_questions": [
                f"Who is responsible for leading the response during {phase}?",
                f"What are the immediate priorities in the {phase} phase?",
                f"How does this affect {industry} regulatory obligations?",
                "What external resources or vendors need to be engaged?",
                "What communication is required internally and externally?",
            ],
            "expected_actions": [
                f"Activate incident response plan",
                f"Assemble the incident response team",
                f"Begin documentation of the {phase.lower()} timeline",
                f"Notify relevant stakeholders",
            ],
        })
    return {
        "overview": f"This tabletop exercise simulates a {template['title']} scenario affecting {client_name}, a {industry} organization.",
        "objectives": [
            "Test incident response plan effectiveness",
            "Identify gaps in communication protocols",
            "Validate decision-making under pressure",
            "Verify regulatory notification procedures",
            "Assess cross-functional coordination",
        ],
        "phases": phases,
        "debrief_questions": [
            "What went well during the exercise?",
            "Where did the team encounter the most friction?",
            "Were escalation paths clear and followed?",
            "What policy or procedure gaps were identified?",
            "What are the top 3 action items from this exercise?",
        ],
    }


async def generate_tabletop_scenario(
    client_id: str,
    scenario_type: str = "ransomware",
    custom_prompt: str = "",
) -> dict:
    """Generate a full tabletop exercise scenario using AI."""
    from db import database as db

    client = await db.get_client(client_id)
    client_name = client.get("name", "the organization") if client else "the organization"
    stack = (client.get("stack_profile") or {}) if client else {}
    industry = client.get("industry") or stack.get("industry") or "Technology"

    template = SCENARIO_TEMPLATES.get(scenario_type, SCENARIO_TEMPLATES["ransomware"])
    hint = custom_prompt or template["prompt_hint"]

    scenario_prompt = (
        f"You are a cybersecurity tabletop exercise facilitator. Generate a detailed tabletop exercise "
        f"for {client_name}, a {industry} organization.\n\n"
        f"Scenario: {hint}\n\n"
        f"Generate a JSON response with this exact structure:\n"
        f'{{"overview": "2-3 sentence scenario overview",\n'
        f'"objectives": ["objective1", "objective2", "objective3", "objective4", "objective5"],\n'
        f'"phases": [\n'
        f'  {{"phase": "Phase Name", "situation": "situation description", "inject": "inject statement", '
        f'"discussion_questions": ["q1","q2","q3","q4","q5"], "expected_actions": ["a1","a2","a3"]}}\n'
        f'],\n'
        f'"debrief_questions": ["q1","q2","q3","q4","q5"]}}\n\n'
        f"Include exactly {len(template['phases'])} phases named: {', '.join(template['phases'])}.\n"
        f"Make it realistic, specific to {industry}, and actionable. Return only valid JSON."
    )

    ai_response = await _ollama_generate(scenario_prompt)
    scenario_data = None

    if ai_response:
        # Extract JSON from response
        try:
            start = ai_response.find("{")
            end = ai_response.rfind("}") + 1
            if start >= 0 and end > start:
                scenario_data = json.loads(ai_response[start:end])
        except Exception:
            pass

    if not scenario_data:
        scenario_data = _fallback_scenario(template, client_name, industry)

    # Save to DB
    title = f"{template['title']} — {client_name} ({datetime.utcnow().strftime('%Y-%m-%d')})"
    saved = await db.create_tabletop(client_id, title, scenario_type, scenario_data)

    return {
        "id": saved["id"],
        "client_id": client_id,
        "client_name": client_name,
        "industry": industry,
        "scenario_type": scenario_type,
        "title": title,
        "scenario": scenario_data,
        "generated_at": saved["created_at"],
        "ai_generated": bool(ai_response),
    }


def export_tabletop_pdf(tabletop: dict) -> bytes:
    """Export a tabletop scenario as PDF using reportlab."""
    try:
        from reportlab.lib.pagesizes import letter
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import inch
        from reportlab.lib import colors
        from reportlab.platypus import (
            SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
        )
        from reportlab.lib.enums import TA_CENTER, TA_LEFT

        buf = io.BytesIO()
        doc = SimpleDocTemplate(buf, pagesize=letter, topMargin=0.75*inch, bottomMargin=0.75*inch)
        styles = getSampleStyleSheet()

        TEAL = colors.HexColor("#00d4aa")
        NAVY = colors.HexColor("#1a2035")
        DARK = colors.HexColor("#0d0f14")

        title_style = ParagraphStyle("title", parent=styles["Heading1"],
                                     fontSize=18, textColor=TEAL, spaceAfter=6)
        h2_style = ParagraphStyle("h2", parent=styles["Heading2"],
                                   fontSize=12, textColor=TEAL, spaceAfter=4, spaceBefore=14)
        h3_style = ParagraphStyle("h3", parent=styles["Heading3"],
                                   fontSize=10, textColor=colors.HexColor("#e8a530"), spaceAfter=4)
        body_style = ParagraphStyle("body", parent=styles["Normal"], fontSize=9, spaceAfter=4)
        bullet_style = ParagraphStyle("bullet", parent=styles["Normal"],
                                       fontSize=9, leftIndent=14, spaceAfter=2, bulletIndent=0)

        story = []
        sc = tabletop.get("scenario", {})

        story.append(Paragraph(tabletop.get("title", "Tabletop Exercise"), title_style))
        story.append(Paragraph(
            f"Client: {tabletop.get('client_name', '')} | Industry: {tabletop.get('industry', '')} | "
            f"Date: {tabletop.get('generated_at', '')[:10]}",
            ParagraphStyle("meta", parent=styles["Normal"], fontSize=8,
                           textColor=colors.grey, spaceAfter=12)
        ))

        story.append(Paragraph("Scenario Overview", h2_style))
        story.append(Paragraph(sc.get("overview", ""), body_style))
        story.append(Spacer(1, 8))

        objs = sc.get("objectives", [])
        if objs:
            story.append(Paragraph("Exercise Objectives", h2_style))
            for obj in objs:
                story.append(Paragraph(f"• {obj}", bullet_style))
            story.append(Spacer(1, 8))

        for phase in sc.get("phases", []):
            story.append(Paragraph(f"Phase: {phase.get('phase', '')}", h2_style))
            story.append(Paragraph(phase.get("situation", ""), body_style))
            inj = phase.get("inject", "")
            if inj:
                story.append(Paragraph(f"<b>INJECT:</b> {inj}", ParagraphStyle(
                    "inject", parent=styles["Normal"], fontSize=9,
                    textColor=colors.HexColor("#e8a530"), spaceAfter=6, spaceBefore=4)))
            story.append(Paragraph("Discussion Questions", h3_style))
            for q in phase.get("discussion_questions", []):
                story.append(Paragraph(f"• {q}", bullet_style))
            story.append(Paragraph("Expected Actions", h3_style))
            for a in phase.get("expected_actions", []):
                story.append(Paragraph(f"→ {a}", bullet_style))
            story.append(Spacer(1, 10))

        debriefs = sc.get("debrief_questions", [])
        if debriefs:
            story.append(PageBreak())
            story.append(Paragraph("Debrief Questions", h2_style))
            for q in debriefs:
                story.append(Paragraph(f"• {q}", bullet_style))

        story.append(Spacer(1, 20))
        story.append(Paragraph(
            "Generated by PhantomFeed Threat Intelligence Platform — CONFIDENTIAL",
            ParagraphStyle("footer", parent=styles["Normal"], fontSize=7,
                           textColor=colors.grey, alignment=TA_CENTER)
        ))

        doc.build(story)
        buf.seek(0)
        return buf.read()

    except ImportError:
        return b""


def export_tabletop_pptx(tabletop: dict) -> bytes:
    """Export a tabletop scenario as PPTX."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor

        NAVY = RGBColor(0x1a, 0x20, 0x35)
        TEAL = RGBColor(0x00, 0xd4, 0xaa)
        WHITE = RGBColor(0xFF, 0xFF, 0xFF)
        AMBER = RGBColor(0xE8, 0xA5, 0x30)
        GRAY = RGBColor(0x64, 0x74, 0x8b)

        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6]
        sc = tabletop.get("scenario", {})

        def add_slide():
            sl = prs.slides.add_slide(blank)
            bg = sl.background.fill
            bg.solid()
            bg.fore_color.rgb = NAVY
            return sl

        def tb(sl, l, t, w, h, text, fs=11, bold=False, color=None):
            box = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = text
            r.font.size = Pt(fs)
            r.font.bold = bold
            r.font.color.rgb = color or WHITE

        def hdr(sl, title, sub=""):
            bar = sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.1))
            bar.fill.solid(); bar.fill.fore_color.rgb = TEAL; bar.line.fill.background()
            tb(sl, 0.2, 0.08, 9.6, 0.6, title, 22, bold=True, color=NAVY)
            if sub:
                tb(sl, 0.2, 0.72, 9.6, 0.35, sub, 9, color=NAVY)

        # Title slide
        sl = add_slide()
        bar = sl.shapes.add_shape(1, Inches(0), Inches(6.5), Inches(10), Inches(1))
        bar.fill.solid(); bar.fill.fore_color.rgb = TEAL; bar.line.fill.background()
        tb(sl, 0.5, 1.5, 9, 0.8, "TABLETOP EXERCISE", 28, bold=True, color=TEAL)
        tb(sl, 0.5, 2.5, 9, 0.6, tabletop.get("title", ""), 16, bold=True)
        tb(sl, 0.5, 3.2, 9, 0.4, f"Client: {tabletop.get('client_name','')} | {tabletop.get('industry','')}", 11, color=GRAY)
        tb(sl, 0.5, 3.7, 9, 0.4, f"CONFIDENTIAL — {tabletop.get('generated_at','')[:10]}", 9, color=GRAY)

        # Overview + objectives
        sl = add_slide()
        hdr(sl, "Scenario Overview & Objectives")
        tb(sl, 0.3, 1.2, 9.5, 0.7, sc.get("overview", ""), 11, color=GRAY)
        y = 2.0
        for obj in (sc.get("objectives") or [])[:5]:
            tb(sl, 0.3, y, 9.5, 0.4, f"• {obj}", 10)
            y += 0.5

        # One slide per phase
        for phase in sc.get("phases", []):
            sl = add_slide()
            hdr(sl, f"Phase: {phase.get('phase', '')}", phase.get("situation", "")[:100])
            inj = phase.get("inject", "")
            if inj:
                tb(sl, 0.3, 1.2, 9.5, 0.4, f"INJECT: {inj[:120]}", 9, color=AMBER)
            tb(sl, 0.3, 1.7, 4.5, 0.3, "Discussion Questions", 10, bold=True, color=TEAL)
            y = 2.05
            for q in (phase.get("discussion_questions") or [])[:4]:
                tb(sl, 0.3, y, 4.5, 0.4, f"• {q[:80]}", 9, color=GRAY)
                y += 0.45
            tb(sl, 5.2, 1.7, 4.5, 0.3, "Expected Actions", 10, bold=True, color=TEAL)
            y = 2.05
            for a in (phase.get("expected_actions") or [])[:4]:
                tb(sl, 5.2, y, 4.5, 0.4, f"→ {a[:80]}", 9, color=GRAY)
                y += 0.45

        # Debrief
        sl = add_slide()
        hdr(sl, "Debrief Questions")
        y = 1.3
        for q in (sc.get("debrief_questions") or [])[:6]:
            tb(sl, 0.3, y, 9.5, 0.5, f"• {q}", 11)
            y += 0.65

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return buf.read()

    except ImportError:
        return b""
